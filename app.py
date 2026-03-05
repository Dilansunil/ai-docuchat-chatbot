import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.llms import Ollama
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
import hashlib
import os
import sqlite3
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime
import pyttsx3
from googletrans import Translator
import speech_recognition as sr
from queue import Queue
import docx
from pptx import Presentation
import pandas as pd

def init_db():
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS users
                 (username TEXT PRIMARY KEY, password TEXT)''')
    conn.commit()
    conn.close()

init_db()

def register_user(username, password):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    try:
        c.execute("INSERT INTO users VALUES (?, ?)", (username, password))
        conn.commit()
        return True
    except sqlite3.IntegrityError:
        return False
    finally:
        conn.close()

def authenticate_user(username, password):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute("SELECT * FROM users WHERE username=? AND password=?", (username, password))
    result = c.fetchone()
    conn.close()
    return result is not None

def init_session():
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
    if "username" not in st.session_state:
        st.session_state.username = None
    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "vectorstore" not in st.session_state:
        st.session_state.vectorstore = None
    if "processed_docs" not in st.session_state:
        st.session_state.processed_docs = []
    if "dark_mode" not in st.session_state:
        st.session_state.dark_mode = False
    if "target_lang" not in st.session_state:
        st.session_state.target_lang = 'en'
    if "voice_enabled" not in st.session_state:
        st.session_state.voice_enabled = False
    if "listening" not in st.session_state:
        st.session_state.listening = False
    if "audio_queue" not in st.session_state:
        st.session_state.audio_queue = Queue()
    if "stop_listening_flag" not in st.session_state:
        st.session_state.stop_listening_flag = False
    if "audio_input" not in st.session_state:
        st.session_state.audio_input = ""
    if "recognizer" not in st.session_state:
        st.session_state.recognizer = sr.Recognizer()
    if "microphone" not in st.session_state:
        st.session_state.microphone = sr.Microphone()
    if "recognizer" not in st.session_state:
        st.session_state.recognizer = sr.Recognizer()
        st.session_state.recognizer.energy_threshold = 4000
        st.session_state.recognizer.dynamic_energy_threshold = True
    if "microphone" not in st.session_state:
        try:
            st.session_state.microphone = sr.Microphone()
            st.session_state.mic_available = True
        except:
            st.session_state.mic_available = False
    if "listening" not in st.session_state:
        st.session_state.listening = False
    if "audio_thread" not in st.session_state:
        st.session_state.audio_thread = None

def extract_pdf_text(pdf_file):
    try:
        return " ".join(page.extract_text() or "" for page in PdfReader(pdf_file).pages)
    except Exception as e:
        st.error(f"Error reading PDF {pdf_file.name}: {str(e)}")
        return ""

def extract_docx_text(docx_file):
    try:
        doc = docx.Document(docx_file)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        st.error(f"Error reading DOCX {docx_file.name}: {str(e)}")
        return ""

def extract_pptx_text(pptx_file):
    try:
        prs = Presentation(pptx_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        st.error(f"Error reading PPTX {pptx_file.name}: {str(e)}")
        return ""

def extract_xlsx_text(xlsx_file):
    try:
        df = pd.read_excel(xlsx_file, sheet_name=None)
        text = []
        for sheet_name, sheet_data in df.items():
            text.append(f"Sheet: {sheet_name}")
            text.append(sheet_data.to_string())
        return "\n".join(text)
    except Exception as e:
        st.error(f"Error reading XLSX {xlsx_file.name}: {str(e)}")
        return ""

def extract_txt_text(txt_file):
    try:
        return txt_file.read().decode("utf-8")
    except Exception as e:
        st.error(f"Error reading TXT {txt_file.name}: {str(e)}")
        return ""

def extract_text(file):
    file_type = file.name.split(".")[-1].lower()
    if file_type == "pdf":
        return extract_pdf_text(file)
    elif file_type == "docx":
        return extract_docx_text(file)
    elif file_type == "pptx":
        return extract_pptx_text(file)
    elif file_type == "xlsx":
        return extract_xlsx_text(file)
    elif file_type == "txt":
        return extract_txt_text(file)
    else:
        st.error(f"Unsupported file type: {file_type}")
        return ""

def process_files_parallel(files):
    with ThreadPoolExecutor(max_workers=4) as executor:
        return " ".join(executor.map(extract_text, files))

def chunk_text(text):
    return RecursiveCharacterTextSplitter(
        chunk_size=768,
        chunk_overlap=100,
        length_function=len
    ).split_text(text)

def get_vectorstore(text_chunks):
    cache_key = hashlib.md5("".join(text_chunks).encode()).hexdigest()
    cache_path = f"vector_cache/{st.session_state.username}_{cache_key}"
    
    if os.path.exists(cache_path):
        try:
            return FAISS.load_local(
                cache_path, 
                OllamaEmbeddings(model="nomic-embed-text"),
                allow_dangerous_deserialization=True
            )
        except Exception as e:
            st.error(f"Error loading cache: {str(e)}")
            os.remove(cache_path)
    
    embeddings = OllamaEmbeddings(model="nomic-embed-text")
    vectorstore = FAISS.from_texts(text_chunks, embedding=embeddings)
    os.makedirs("vector_cache", exist_ok=True)
    vectorstore.save_local(cache_path)
    return vectorstore

def init_voice_engine():
    try:
        engine = pyttsx3.init()
        voices = engine.getProperty('voices')
        engine.setProperty('voice', voices[0].id)
        engine.setProperty('rate', 150)
        return engine
    except Exception as e:
        st.error(f"Voice engine error: {str(e)}")
        return None

def translate_text(text, dest_lang='en'):
    try:
        translator = Translator()
        return translator.translate(text, dest=dest_lang).text
    except Exception as e:
        st.error(f"Translation error: {str(e)}")
        return text

def login_page():
    st.title("Login to DocuChat")
    username = st.text_input("Username")
    password = st.text_input("Password", type="password")
    
    if st.button("Login"):
        if authenticate_user(username, password):
            st.session_state.authenticated = True
            st.session_state.username = username
            st.rerun()
        else:
            st.error("Invalid username or password")
    
    if st.button("Create New Account"):
        st.session_state.show_register = True
        st.rerun()

def register_page():
    st.title("Create Account")
    username = st.text_input("Username")
    password = st.text_input("Password", type="password")
    confirm_password = st.text_input("Confirm Password", type="password")
    
    if st.button("Register"):
        if password != confirm_password:
            st.error("Passwords don't match")
        elif len(username) < 3:
            st.error("Username must be at least 3 characters")
        elif len(password) < 6:
            st.error("Password must be at least 6 characters")
        elif register_user(username, password):
            st.success("Registration successful! Please login.")
            st.session_state.show_register = False
            st.rerun()
        else:
            st.error("Username already exists")
    
    if st.button("Back to Login"):
        st.session_state.show_register = False
        st.rerun()

def handle_query(query):
    if not st.session_state.conversation:
        st.warning("Please process documents first")
        return
    
    with st.spinner("Thinking..."):
        try:
            response = st.session_state.conversation({"question": query})
            st.session_state.chat_history = response["chat_history"]
            
            if st.session_state.voice_enabled:
                engine = init_voice_engine()
                if engine:
                    translated_response = translate_text(response['answer'], st.session_state.target_lang)
                    engine.say(translated_response)
                    engine.runAndWait()
            
            st.rerun()
        except Exception as e:
            st.error(f"Error generating response: {str(e)}")

def check_microphone():
    try:
        with st.session_state.microphone as source:
            st.session_state.recognizer.adjust_for_ambient_noise(source, duration=1)
        return True
    except Exception as e:
        st.error(f"Microphone error: {str(e)}")
        return False

def chat_interface():
    if st.session_state.dark_mode:
        st.markdown("""
        <style>
            :root {
                --bg-color: #1e1e1e;
                --text-color: #ffffff;
                --bot-bg: #333333;
                --user-bg: #4a6baf;
            }
            .main { background-color: var(--bg-color); color: var(--text-color); }
            .sidebar .sidebar-content { background-color: #2d2d2d; }
            .stTextInput input { color: var(--text-color); background-color: #333; }
            .bot-message { background-color: var(--bot-bg) !important; color: var(--text-color) !important; }
            .user-message { background-color: var(--user-bg) !important; color: white !important; }
        </style>
        """, unsafe_allow_html=True)
    
    with st.sidebar:
        st.title("DocuChat Settings")
        
        st.session_state.target_lang = st.selectbox(
            "Select Language",
            ['en', 'es', 'fr', 'de', 'it', 'pt', 'ru', 'zh', 'ja'],
            index=0
        )
        
        st.session_state.voice_enabled = st.checkbox("Enable Voice Output", value=False)
        
        st.markdown("---")
        st.subheader("Document Processing")
        uploaded_files = st.file_uploader(
            "Upload Documents", 
            accept_multiple_files=True,
            type=["pdf", "docx", "pptx", "xlsx", "txt"]
        )
        
        if st.button("Process Documents"):
            if uploaded_files:
                with st.status("Processing documents...", expanded=True) as status:
                    st.write("Extracting text...")
                    text = process_files_parallel(uploaded_files)
                    
                    if not text.strip():
                        st.error("No text could be extracted")
                        return
                    
                    st.write("Creating chunks...")
                    chunks = chunk_text(text)
                    
                    st.write("Building vector store...")
                    st.session_state.vectorstore = get_vectorstore(chunks)
                    st.session_state.processed_docs = [file.name for file in uploaded_files]
                    
                    st.write("Initializing chat...")
                    st.session_state.conversation = ConversationalRetrievalChain.from_llm(
                        llm=Ollama(model="llama2:7b-chat"),
                        retriever=st.session_state.vectorstore.as_retriever(),
                        memory=ConversationBufferMemory(
                            memory_key='chat_history',
                            return_messages=True
                        )
                    )
            else:
                st.warning("Please upload documents first")
        
        if st.button("Logout"):
            st.session_state.authenticated = False
            st.session_state.username = None
            st.rerun()
    
    st.title(f"Welcome, {st.session_state.username}!")
    st.subheader("DocuChat - Train Your AI on Documents")
    
    if st.session_state.processed_docs:
        st.markdown("**Processed Documents:**")
        for doc in st.session_state.processed_docs:
            st.markdown(f"- {doc}")
    
    if st.session_state.vectorstore:
        st.markdown("**Suggested Questions:**")
        cols = st.columns(3)
        questions = [
            "What is this document about?",
            "Summarize the key points",
            "What are the main conclusions?"
        ]
        for i, q in enumerate(questions):
            if cols[i%3].button(q):
                handle_query(q)

    chat_container = st.container()
    with chat_container:
        for i, msg in enumerate(st.session_state.chat_history):
            translated = translate_text(msg.content, st.session_state.target_lang)
            
            if i % 2 == 0:  
                st.markdown(f"""
                <div style='background-color: #4a6baf; color: white; 
                            border-radius: 15px; padding: 12px; margin: 10px 0;
                            max-width: 80%; margin-left: auto;'>
                    {translated}
                    <div style='font-size: 0.7rem; text-align: right; margin-top: 5px;'>
                        {datetime.now().strftime('%H:%M')}
                    </div>
                </div>
                """, unsafe_allow_html=True)
            else: 
                st.markdown(f"""
                <div style='background-color: {'#333' if st.session_state.dark_mode else '#f0f0f0'}; 
                            color: {'white' if st.session_state.dark_mode else '#333'}; 
                            border-radius: 15px; padding: 12px; margin: 10px 0;
                            max-width: 80%; margin-right: auto;'>
                    {translated}
                    <div style='font-size: 0.7rem; text-align: right; margin-top: 5px;
                                color: {'#aaa' if st.session_state.dark_mode else '#666'};'>
                        {datetime.now().strftime('%H:%M')}
                    </div>
                </div>
                """, unsafe_allow_html=True)
    
    if user_input := st.chat_input(f"Ask about your documents (in {st.session_state.target_lang})..."):
        handle_query(user_input)

def main():
    init_session()
    
    if not st.session_state.authenticated:
        if st.session_state.get('show_register'):
            register_page()
        else:
            login_page()
    else:
        chat_interface()

if __name__ == "__main__":
    main()